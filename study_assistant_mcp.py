from pathlib import Path
from typing import Any, Dict, List, TypedDict

from mcp.server.fastmcp import FastMCP

server = FastMCP("study-assistant-mcp")

BASE_DIR = Path(__file__).parent
DATA_DIR = BASE_DIR / "data"
DATA_DIR.mkdir(exist_ok=True)

STUDY_DIR = DATA_DIR / "study"
STUDY_DIR.mkdir(exist_ok=True)

PYQ_DIR = DATA_DIR / "pyq"
PYQ_DIR.mkdir(exist_ok=True)

#  Build a safe path for a file in data/study.
def _safe_study_file(name: str) -> Path: 
    if "/" in name or "\\" in name:
        raise ValueError("Study filename must not contain path separators.")
    return STUDY_DIR / name

# Build a safe path for a file in data/pyq.
# PYQ files are stored as .txt by default.
def _safe_pyq_file(name: str) -> Path:
    if "/" in name or "\\" in name:
        raise ValueError("PYQ name must not contain path separators.")
    if not name.endswith(".txt"):
        name = name + ".txt"
    return PYQ_DIR / name

@server.tool()
def list_study_files(extensions: List[str] | None = None) -> List[str]:
    if extensions is None or len(extensions) == 0:
        exts: List[str] = []
    else:
        exts = [e.lower().lstrip(".") for e in extensions]

    files: List[str] = []
    for p in STUDY_DIR.iterdir():
        if not p.is_file():
            continue
        if exts:
            if p.suffix.lower().lstrip(".") not in exts:
                continue
        files.append(p.name)
    return files

# Extract plain text from a PDF in 'data/study'.
#     - max_pages: limit to first N pages
#     - max_chars: truncate result to this many characters.
@server.tool()
def extract_pdf_text(
    filename: str,
    max_pages: int = 30,
    max_chars: int = 8000,
) -> str:
    from pypdf import PdfReader

    file_path = _safe_study_file(filename)
    if not file_path.exists():
        raise FileNotFoundError(
            f"Study file '{filename}' does not exist in data/study"
        )

    reader = PdfReader(str(file_path))
    texts: List[str] = []
    page_count = min(len(reader.pages), max_pages)

    for i in range(page_count):
        page = reader.pages[i]
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt.strip():
            texts.append(f"[Page {i+1}]\n{txt.strip()}")

    full_text = "\n\n".join(texts).strip()
    if not full_text:
        return "No extractable text found in this PDF."

    if len(full_text) > max_chars:
        return full_text[: max_chars - 3] + "..."

    return full_text

# Extract text from a PPTX file in 'data/study'.
    # Concatenates text from all slides, truncated to max_chars.
@server.tool()
def extract_pptx_text(
    filename: str,
    max_chars: int = 8000,
) -> str:
    from pptx import Presentation

    file_path = _safe_study_file(filename)
    if not file_path.exists():
        raise FileNotFoundError(
            f"Study file '{filename}' does not exist in data/study"
        )

    prs = Presentation(str(file_path))
    parts: List[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))

    content = "\n\n".join(parts).strip()
    if not content:
        return "No extractable text found in this PPTX."

    if len(content) > max_chars:
        return content[: max_chars - 3] + "..."

    return content


#   PYQ TOOLS
@server.tool()
# Save a set of previous year questions (PYQs) as a text file
    # in the 'data/pyq' folder. Stored as <name>.txt if no extension is given.
def save_pyq_set(name: str, content: str) -> str:
    file_path = _safe_pyq_file(name)
    with file_path.open("w", encoding="utf-8") as f:
        f.write(content.strip() + "\n")
    return f"Saved PYQs as {file_path.name}."


@server.tool()
# List available PYQ sets (filenames) in 'data/pyq'.
def list_pyq_sets() -> List[str]:
    names: List[str] = []
    for p in PYQ_DIR.iterdir():
        if p.is_file():
            names.append(p.name)
    return names


@server.tool()
#  Read the content of a PYQ set from 'data/pyq'.
def read_pyq_set(name: str) -> str:
    file_path = _safe_pyq_file(name)
    if not file_path.exists():
        raise FileNotFoundError(
            f"PYQ set '{file_path.name}' does not exist in data/pyq"
        )

    with file_path.open("r", encoding="utf-8") as f:
        return f.read()

#   BUNDLE FOR EXAM QUESTION GENERATION
class StudyBundle(TypedDict):
    study_sources: List[str]
    pyq_sources: List[str]
    text: str

#  Internal helper used by both the MCP tool and the API-based question generator.
def _build_study_bundle_core(
    study_files: List[str],
    pyq_sets: List[str],
    max_chars: int = 12000,
) -> StudyBundle:
    parts: List[str] = []
    used_study: List[str] = []
    used_pyq: List[str] = []

    # Study files: auto-detect pdf/pptx by extension
    for fname in study_files:
        lower = fname.lower()
        if lower.endswith(".pdf"):
            snippet = extract_pdf_text(
                fname,
                max_pages=40,
                max_chars=max_chars // 2,
            )
        elif lower.endswith(".pptx"):
            snippet = extract_pptx_text(
                fname,
                max_chars=max_chars // 2,
            )
        else:
            # Fallback: treat as plain text in study/ folder
            file_path = _safe_study_file(fname)
            if not file_path.exists():
                continue
            snippet = file_path.read_text(encoding="utf-8")

        parts.append(f"=== Study File: {fname} ===\n{snippet}\n")
        used_study.append(fname)

    # PYQ sets
    for name in pyq_sets:
        text = read_pyq_set(name)
        label = name if name.endswith(".txt") else name + ".txt"
        parts.append(f"=== PYQs: {label} ===\n{text.strip()}\n")
        used_pyq.append(label)

    combined = "\n\n".join(parts).strip()
    if len(combined) > max_chars:
        combined = combined[: max_chars - 3] + "..."

    return {
        "study_sources": used_study,
        "pyq_sources": used_pyq,
        "text": combined,
    }

# Build a single text bundle from study files + PYQs.
@server.tool()
def build_study_bundle(
    study_files: List[str],
    pyq_sets: List[str],
    max_chars: int = 12000,
) -> StudyBundle:
    return _build_study_bundle_core(study_files, pyq_sets, max_chars)

@server.tool()
def generate_probable_questions(
    study_files: List[str],
    pyq_sets: List[str],
    max_questions: int = 20,
    max_chars: int = 12000,
    model: str = "llama3.2:1b",
) -> str:
    # Use a local Ollama model to generate probable exam questions.
    import ollama

    bundle = _build_study_bundle_core(
        study_files=study_files,
        pyq_sets=pyq_sets,
        max_chars=max_chars,
    )

    context_text = bundle["text"]
    if not context_text:
        return "No content available from the selected study files and PYQs."

    system_msg = (
        "You are an expert university instructor helping a student prepare for an exam.\n"
        "You are given a combination of course materials (slides/notes) and previous year questions.\n"
        "Your task is to generate likely exam questions that match the style and difficulty of the course.\n"
        "Only output the questions, numbered, without answers.\n"
    )

    user_msg = (
        f"Generate {max_questions} answers based on the following content, and the answers should be just the option number in the format:\n"
        "Answer 1. - (correct answer)"
        "----- COMBINED STUDY MATERIAL + PYQs -----\n"
        f"{context_text}\n"
        "------------------------------------------\n"
    )

    response = ollama.chat(
        model=model,
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_msg},
        ],
    )

    content = response["message"]["content"]
    return content.strip()



# Main
def main() -> None:
    server.run(transport="stdio")


if __name__ == "__main__":
    main()
