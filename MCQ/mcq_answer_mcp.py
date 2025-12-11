from pathlib import Path
from typing import List, Dict

from mcp.server.fastmcp import FastMCP

import re  # for MCQ splitting and validation

server = FastMCP("mcq-answer-mcp")

BASE_DIR = Path(__file__).parent
DATA_DIR = BASE_DIR / "data"
DATA_DIR.mkdir(exist_ok=True)

PYQ_RAW_DIR = DATA_DIR / "pyq_raw"
PYQ_RAW_DIR.mkdir(exist_ok=True)

PYQ_TXT_DIR = DATA_DIR / "pyq_txt"
PYQ_TXT_DIR.mkdir(exist_ok=True)

ANSWERS_DIR = DATA_DIR / "answers"
ANSWERS_DIR.mkdir(exist_ok=True)


def _safe_raw_file(name: str) -> Path:
    if "/" in name or "\\" in name:
        raise ValueError("Filename must not contain path separators.")
    return PYQ_RAW_DIR / name


def _safe_txt_file(name: str) -> Path:
    if "/" in name or "\\" in name:
        raise ValueError("Filename must not contain path separators.")
    if not name.endswith(".txt"):
        name = name + ".txt"
    return PYQ_TXT_DIR / name


def _safe_answer_file(name: str) -> Path:
    if "/" in name or "\\" in name:
        raise ValueError("Filename must not contain path separators.")
    if not name.endswith(".txt"):
        name = name + ".txt"
    return ANSWERS_DIR / name


def _extract_section_a(full_text: str) -> str:
    """
    Extract (or approximate) Section A from the exam text.

    Strategy:
    - Use the *last* occurrence of "Section A" as the start (this is the real
      Section A header above the MCQs, not the earlier "Section Acarries 40%..." line).
    - Then cut at the first "Question B" or "Section B" after that, if present.
    - If we can't find anything reasonable, fall back to the full text.
    """
    text = full_text

    # Start from the LAST "Section A" – this should be the header above MCQs.
    start_idx = text.rfind("Section A")
    if start_idx == -1:
        # No explicit Section A marker – just return everything.
        return full_text.strip()

    text = text[start_idx:]

    # Now try to find the end of Section A.
    # Prefer an explicit "Question B" (start of Section B questions), otherwise "Section B".
    qb_idx = text.find("Question B")
    if qb_idx != -1:
        text = text[:qb_idx]
    else:
        sb_idx = text.find("Section B")
        if sb_idx != -1:
            text = text[:sb_idx]

    return text.strip()



import re

def _split_mcq_blocks(section_a_text: str) -> List[tuple[int, str]]:
    """
    Split Section A text into individual MCQ blocks.

    Designed for patterns like:
      'Question A 1What is supervised learning?'
      'Question A 2Which one of the following is true?'

    We allow missing / weird spaces after 'A' and after the number.
    Returns list of (question_number, question_block).
    """
    pattern = re.compile(
        r"(Question\s*A\s*(\d+).*?)(?=Question\s*A\s*\d+|$)",
        re.DOTALL,
    )
    blocks: List[tuple[int, str]] = []
    for match in pattern.finditer(section_a_text):
        block_text = match.group(1).strip()
        q_num = int(match.group(2))
        blocks.append((q_num, block_text))
    return blocks


@server.tool()
def list_pyq_raw_files(extensions: List[str] | None = None) -> List[str]:
    # List files in data/pyq_raw, optionally filtered by extension.
    if not extensions:
        exts: List[str] = []
    else:
        exts = [e.lower().lstrip(".") for e in extensions]

    files: List[str] = []
    for p in PYQ_RAW_DIR.iterdir():
        if not p.is_file():
            continue
        if exts:
            if p.suffix.lower().lstrip(".") not in exts:
                continue
        files.append(p.name)
    return files


@server.tool()
def extract_pyq_pdf_to_txt(filename: str, output_name: str) -> str:
    # Extract text from a PDF in data/pyq_raw into a txt file in data/pyq_txt.
    from pypdf import PdfReader

    raw_path = _safe_raw_file(filename)
    if not raw_path.exists():
        raise FileNotFoundError(f"File '{filename}' does not exist in data/pyq_raw")

    reader = PdfReader(str(raw_path))
    texts: List[str] = []

    for i, page in enumerate(reader.pages, start=6):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt.strip():
            texts.append(f"[Page {i}]\n{txt.strip()}")

    full_text = "\n\n".join(texts).strip()
    if not full_text:
        full_text = "No extractable text found in this PDF."

    out_path = _safe_txt_file(output_name)
    with out_path.open("w", encoding="utf-8") as f:
        f.write(full_text + "\n")

    return f"Extracted text from {filename} to {out_path.name}."


@server.tool()
def extract_pyq_pptx_to_txt(filename: str, output_name: str) -> str:
    # Extract text from a PPTX in data/pyq_raw into a txt file in data/pyq_txt.
    from pptx import Presentation

    raw_path = _safe_raw_file(filename)
    if not raw_path.exists():
        raise FileNotFoundError(f"File '{filename}' does not exist in data/pyq_raw")

    prs = Presentation(str(raw_path))
    parts: List[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    slide_texts.append(t)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

    full_text = "\n\n".join(parts).strip()
    if not full_text:
        full_text = "No extractable text found in this PPTX."

    out_path = _safe_txt_file(output_name)
    with out_path.open("w", encoding="utf-8") as f:
        f.write(full_text + "\n")

    return f"Extracted text from {filename} to {out_path.name}."


@server.tool()
def list_pyq_txt_files() -> List[str]:
    # List .txt PYQ files in data/pyq_txt.
    files: List[str] = []
    for p in PYQ_TXT_DIR.iterdir():
        if p.is_file() and p.suffix.lower() == ".txt":
            files.append(p.name)
    return files

@server.tool()
def preview_pyq_txt(pyq_txt_filename: str, max_chars: int = 2000) -> str:
    """
    Return the first part of a txt file from data/pyq_txt for debugging.
    """
    txt_path = _safe_txt_file(pyq_txt_filename)
    if not txt_path.exists():
        raise FileNotFoundError(
            f"PYQ txt file '{txt_path.name}' does not exist in data/pyq_txt"
        )

    with txt_path.open("r", encoding="utf-8") as f:
        content = f.read()

    snippet = content[:max_chars]
    return f"Preview of {txt_path.name} (first {len(snippet)} chars):\n\n{snippet}"


@server.tool()
def debug_preview_file(pyq_txt_filename: str, max_chars: int = 3000) -> str:
    """
    Preview the first 3000 characters of a txt file in data/pyq_txt.
    Helps verify if the extracted text actually contains Section A MCQs.
    """
    txt_path = _safe_txt_file(pyq_txt_filename)
    if not txt_path.exists():
        raise FileNotFoundError(f"File '{txt_path.name}' does not exist in data/pyq_txt")

    with txt_path.open("r", encoding="utf-8") as f:
        content = f.read()

    snippet = content[:max_chars]
    return f"Preview of {txt_path.name}:\n\n{snippet}"


@server.tool()
def answer_mcq_file(
    pyq_txt_filename: str,
    model: str = "llama3.1:8b",
    output_filename: str | None = None,
) -> str:
    """
    Generate MCQ answers (only options) from a txt file using a local Ollama model.

    - Reads the txt file from data/pyq_txt.
    - Narrows to Section A (MCQs) when possible.
    - Sends the MCQ text to the model with strong instructions.
    - Cleans the model output into lines like "1: A".
    - Writes the result into data/answers.
    """
    import ollama
    import re

    txt_path = _safe_txt_file(pyq_txt_filename)
    if not txt_path.exists():
        raise FileNotFoundError(
            f"PYQ txt file '{txt_path.name}' does not exist in data/pyq_txt"
        )

    with txt_path.open("r", encoding="utf-8") as f:
        full_text = f.read().strip()

    if not full_text:
        return "PYQ txt file is empty."

    # Try to restrict to Section A only (MCQs).
    section_a = _extract_section_a(full_text)
    exam_text = section_a if section_a else full_text

    system_msg = (
        "You are an expert in machine learning exams, solving multiple-choice questions.\n"
        "You will receive the OCR text of a university exam paper.\n"
        "Section A contains exactly 20 multiple-choice questions labelled "
        "'Question A 1', 'Question A 2', ..., 'Question A 20', each with answer "
        "options A, B, C, and D.\n\n"
        "Your job:\n"
        "- Carefully read each question and its four options.\n"
        "- For EACH of the 20 questions (A1–A20), decide which single option "
        "is most correct based on the content of the question and options.\n"
        "- You must base your answers ONLY on the question text and options provided.\n\n"
        "Output format:\n"
        "- Output EXACTLY 20 lines.\n"
        "- Each line must be: '<number>: <LETTER>' where <number> is from 1 to 20,\n"
        "  and <LETTER> is exactly one of A, B, C, or D.\n"
        "  Example:\n"
        "  1: B\n"
        "  2: C\n"
        "  3: D\n"
        "- The number corresponds to the question index (Question A 1 -> '1', ..., "
        "Question A 20 -> '20').\n\n"
        "Very important constraints:\n"
        "- You must NOT output simple patterns like '1: A, 2: B, 3: C, 4: D, 5: A, ...'.\n"
        "- You must think about the meaning of each question and its options.\n"
        "- First, reason internally step by step about each question, but DO NOT "
        "show your reasoning in the final output.\n"
        "- After you have thought about all 20 questions, output ONLY the 20 lines "
        "of answers in the required format.\n"
    )

    user_msg = (
        "Below is the exam text including Section A (MCQ questions and options).\n"
        "Do NOT ask me to provide the questions again; they are already included.\n"
        "Your final output must be exactly 20 lines with answers from 1 to 20.\n\n"
        "----- EXAM TEXT BEGIN -----\n"
        f"{exam_text}\n"
        "----- EXAM TEXT END -----\n"
    )

    response = ollama.chat(
        model=model,
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_msg},
        ],
        # Make it deterministic & a bit more serious about following instructions.
        options={"temperature": 0},
    )

    raw_answers = response["message"]["content"].strip()
    if not raw_answers:
        return "Model did not return any answers."

    # Try to clean into "number: LETTER" lines
    cleaned_lines: List[str] = []
    for line in raw_answers.splitlines():
        m = re.match(
            r"\s*(\d+)\s*[:.\-]?\s*([ABCD])\b",
            line.strip(),
            re.IGNORECASE,
        )
        if m:
            qn = m.group(1)
            letter = m.group(2).upper()
            cleaned_lines.append(f"{qn}: {letter}")

    # If cleaning produced something, prefer it
    if cleaned_lines:
        answers_text = "\n".join(cleaned_lines).strip()
    else:
        # Fallback: at least write whatever the model said, so file is never empty
        answers_text = raw_answers

    # Decide output filename
    if output_filename is None or output_filename.strip() == "":
        base_name = txt_path.stem
        output_filename = f"answers_for_{base_name}.txt"

    ans_path = _safe_answer_file(output_filename)
    with ans_path.open("w", encoding="utf-8") as f:
        f.write(answers_text + "\n")

    return f"Answers written to {ans_path.name}.\n\n{answers_text}"

def main() -> None:
    server.run(transport="stdio")


if __name__ == "__main__":
    main()
